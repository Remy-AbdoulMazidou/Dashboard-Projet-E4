"""
Génère :
  docs/slides_presentation.pptx
  docs/discours_oral_remy.docx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from docx import Document
from docx.shared import Pt as DPt, RGBColor as DRGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH

os.makedirs("docs", exist_ok=True)

# ══════════════════════════════════════════════════════════════════════════════
# COULEURS
# ══════════════════════════════════════════════════════════════════════════════
BG      = RGBColor(0x0F, 0x11, 0x17)   # fond sombre
CARD    = RGBColor(0x1A, 0x1D, 0x2E)   # card
BLUE    = RGBColor(0x2E, 0x86, 0xAB)   # accent principal
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREY    = RGBColor(0xAA, 0xAA, 0xAA)
ACCENT  = RGBColor(0xF2, 0x42, 0x36)   # rouge accent

W = Inches(13.33)
H = Inches(7.5)

# ══════════════════════════════════════════════════════════════════════════════
# HELPERS PPTX
# ══════════════════════════════════════════════════════════════════════════════
def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, x, y, w, h, text="", size=18, bold=False,
        color=WHITE, align=PP_ALIGN.LEFT, bg_color=None, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if bg_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg_color
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox

def bullet_box(slide, x, y, w, h, items, size=16, color=WHITE, title=None):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = title
        run.font.size = Pt(size + 1)
        run.font.bold = True
        run.font.color.rgb = BLUE
        run.font.name = "Calibri"
    for item in items:
        p = tf.paragraphs[0] if (first and not title) else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = "  •  " + item
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Calibri"

def accent_bar(slide, y=0.55, h=0.06):
    bar = slide.shapes.add_shape(1, 0, Inches(y), W, Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()

def card_rect(slide, x, y, w, h, color=CARD):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = BLUE
    return shape

# ══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════════
def build_pptx():
    prs = new_prs()

    # ── SLIDE 1 : TITRE ───────────────────────────────────────────────────────
    s = blank_slide(prs)
    bg(s)
    accent_bar(s, y=0.55, h=0.07)
    box(s, 0.6, 0.9, 12, 1.2,
        "Dashboard analytique — Visualisation des données microstructurales",
        size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    box(s, 0.6, 2.2, 12, 0.6,
        "Rémy ABDOUL MAZIDOU  ·  Projet E4 DSIA  ·  ESIEE Paris  ·  2024-2025",
        size=15, color=GREY, align=PP_ALIGN.LEFT)
    box(s, 0.6, 3.1, 12, 0.5,
        "Soutenance de mi-projet",
        size=14, color=BLUE, align=PP_ALIGN.LEFT)

    # ── SLIDE 2 : MA CONTRIBUTION ─────────────────────────────────────────────
    s = blank_slide(prs)
    bg(s)
    accent_bar(s, y=0.55, h=0.05)
    box(s, 0.6, 0.7, 12, 0.7, "Ma contribution", size=24, bold=True, color=WHITE)

    card_rect(s, 0.6, 1.6, 12, 2.2)
    bullet_box(s, 0.9, 1.7, 11.5, 2.0, [
        "Développement d'un dashboard web interactif en Python",
        "Objectif : rendre les données du projet exploitables sans écrire de code",
        "Visualisation des fichiers CSV produits par l'algorithme d'analyse d'images",
        "Pour l'instant : données simulées, en attente des vraies données",
    ], size=17)

    box(s, 0.6, 4.1, 6, 0.5, "Stack technique", size=14, bold=True, color=BLUE)
    bullet_box(s, 0.6, 4.6, 5.8, 2.0, [
        "Python 3.11 + Plotly Dash",
        "Pandas / NumPy",
        "Dash Bootstrap Components",
        "VS Code",
    ], size=15)
    box(s, 7.0, 4.1, 6, 0.5, "Architecture", size=14, bold=True, color=BLUE)
    bullet_box(s, 7.0, 4.6, 5.8, 2.0, [
        "Code modulaire (un fichier par onglet)",
        "4 onglets indépendants",
        "Filtres interactifs par matériau",
        "Thème sombre, graphiques uniformes",
    ], size=15)

    # ── SLIDE 3 : POURQUOI UN DASHBOARD ──────────────────────────────────────
    s = blank_slide(prs)
    bg(s)
    accent_bar(s, y=0.55, h=0.05)
    box(s, 0.6, 0.7, 12, 0.7, "Pourquoi un dashboard analytique ?", size=24, bold=True, color=WHITE)

    cards = [
        ("Centralisation", "Toutes les données du projet en un seul endroit"),
        ("Accessibilité", "Utilisable sans compétences en programmation"),
        ("Exploration", "Identifier des tendances avant toute modélisation"),
        ("Communication", "Présenter les résultats clairement aux encadrants"),
    ]
    for i, (titre, desc) in enumerate(cards):
        cx = 0.4 + i * 3.2
        card_rect(s, cx, 1.7, 2.9, 2.3)
        box(s, cx + 0.15, 1.85, 2.6, 0.55, titre, size=15, bold=True, color=BLUE)
        box(s, cx + 0.15, 2.45, 2.6, 1.2, desc, size=13, color=WHITE)

    box(s, 0.6, 4.3, 12, 1.5,
        "Dans ce projet, la visualisation est une étape analytique à part entière :\n"
        "elle permet de relier paramètres microstructuraux et propriétés acoustiques,\n"
        "et d'orienter les choix pour la modélisation.",
        size=14, color=GREY)

    # ── SLIDE 4 : PRÉSENTATION DU DASHBOARD ───────────────────────────────────
    s = blank_slide(prs)
    bg(s)
    accent_bar(s, y=0.55, h=0.05)
    box(s, 0.6, 0.7, 12, 0.7, "Le dashboard — 4 onglets", size=24, bold=True, color=WHITE)

    onglets = [
        ("Vue d'ensemble", "KPIs + tableau comparatif\nde tous les matériaux"),
        ("Morphologie des fibres", "Distribution des diamètres\nBoxplot + courbe de densité"),
        ("Propriétés acoustiques", "Courbes d'absorption\nPorosité vs résistivité"),
        ("Comparaison", "Barres par fréquence\nDiamètre vs absorption à 1 kHz"),
    ]
    for i, (nom, desc) in enumerate(onglets):
        cx = 0.4 + i * 3.2
        card_rect(s, cx, 1.7, 2.9, 3.2)
        box(s, cx + 0.1, 1.85, 2.7, 0.7, nom, size=14, bold=True, color=BLUE)
        box(s, cx + 0.1, 2.6, 2.7, 2.0, desc, size=13, color=WHITE)

    box(s, 0.6, 5.2, 12, 0.5,
        "[ Insérer ici une capture d'écran du dashboard ]",
        size=12, color=GREY, align=PP_ALIGN.CENTER)

    # ── SLIDE 5 : DÉMO ────────────────────────────────────────────────────────
    s = blank_slide(prs)
    bg(s)
    accent_bar(s, y=0.55, h=0.05)
    box(s, 0.6, 0.7, 12, 0.7, "Démonstration", size=24, bold=True, color=WHITE)

    card_rect(s, 0.6, 1.6, 12, 4.5)
    box(s, 0.6, 1.6, 12, 4.5,
        "[ Capture d'écran ou démo en direct du dashboard ]",
        size=18, color=GREY, align=PP_ALIGN.CENTER)

    # ── SLIDE 6 : ET APRÈS ────────────────────────────────────────────────────
    s = blank_slide(prs)
    bg(s)
    accent_bar(s, y=0.55, h=0.05)
    box(s, 0.6, 0.7, 12, 0.7, "Et après ?", size=24, bold=True, color=WHITE)

    card_rect(s, 0.6, 1.6, 12, 2.5)
    bullet_box(s, 0.9, 1.75, 11.2, 2.2, [
        "Intégrer les vraies données dès qu'elles seront disponibles",
        "Vérifier la cohérence des graphiques avec des valeurs réelles",
        "Ajouter de nouveaux graphiques selon les besoins du projet",
        "L'outil est conçu pour évoluer facilement",
    ], size=17)

    box(s, 0.6, 4.4, 12, 0.6,
        "Le dashboard est fonctionnel et prêt à recevoir les vraies données.",
        size=15, color=BLUE, bold=True)

    out = "docs/slides_presentation.pptx"
    prs.save(out)
    print(f"PPTX créé : {out}")


# ══════════════════════════════════════════════════════════════════════════════
# DISCOURS ORAL
# ══════════════════════════════════════════════════════════════════════════════
def build_discours():
    doc = Document()

    def h1(text):
        p = doc.add_heading(text, level=1)
        for r in p.runs:
            r.font.name = "Calibri"
            r.font.size = DPt(15)
            r.font.color.rgb = DRGBColor(0x1F, 0x49, 0x7D)

    def h2(text):
        p = doc.add_heading(text, level=2)
        for r in p.runs:
            r.font.name = "Calibri"
            r.font.size = DPt(12)
            r.font.color.rgb = DRGBColor(0x2E, 0x74, 0xB5)

    def body(text):
        p = doc.add_paragraph(text)
        p.paragraph_format.space_after = DPt(6)
        for r in p.runs:
            r.font.name = "Calibri"
            r.font.size = DPt(11)

    def note(text):
        p = doc.add_paragraph(text)
        p.paragraph_format.space_after = DPt(4)
        for r in p.runs:
            r.font.name = "Calibri"
            r.font.size = DPt(10)
            r.font.italic = True
            r.font.color.rgb = DRGBColor(0x70, 0x70, 0x70)

    def graph(titre, onglet, description):
        h2(f"{titre}  —  onglet « {onglet} »")
        body(description)
        doc.add_paragraph()

    # ── EN-TÊTE ───────────────────────────────────────────────────────────────
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("Discours oral — Rémy ABDOUL MAZIDOU")
    r.font.name = "Calibri"
    r.font.size = DPt(16)
    r.font.bold = True

    note("Durée cible : 3 minutes  ·  6 slides  ·  Lire à voix haute de façon naturelle")
    doc.add_paragraph()

    # ── SLIDE 1 ───────────────────────────────────────────────────────────────
    h1("Slide 1 — Titre")
    body(
        "Bonjour à tous. Ma partie dans ce projet, c'est le dashboard analytique. "
        "Concrètement, j'ai développé une application web qui permet de visualiser toutes "
        "les données produites par le projet de façon interactive."
    )

    # ── SLIDE 2 ───────────────────────────────────────────────────────────────
    h1("Slide 2 — Ma contribution")
    body(
        "Ce que j'ai développé, c'est une interface web en Python avec le framework Plotly Dash. "
        "L'idée, c'est que dès que les données seront prêtes — les fichiers CSV produits par "
        "l'algorithme d'analyse d'images — on peut les charger dans le dashboard et les explorer "
        "immédiatement, sans avoir à écrire de code."
    )
    body(
        "En attendant, j'ai généré des données simulées pour pouvoir développer et tester "
        "l'outil. L'architecture est modulaire : chaque onglet correspond à un fichier Python "
        "indépendant, ce qui facilite les ajouts futurs."
    )

    # ── SLIDE 3 ───────────────────────────────────────────────────────────────
    h1("Slide 3 — Pourquoi un dashboard ?")
    body(
        "Je vais prendre un moment pour expliquer pourquoi un dashboard est vraiment utile "
        "dans ce projet, parce que ce n'est pas forcément évident au premier abord."
    )
    body(
        "Ce projet produit beaucoup de données hétérogènes : des fichiers CSV avec des milliers "
        "de lignes décrivant les fibres, leurs diamètres, leurs orientations, et aussi des "
        "mesures acoustiques. Ces données brutes, ouvertes dans un tableur, ne disent rien. "
        "On ne voit pas les tendances, on ne compare pas les matériaux facilement."
    )
    body(
        "Le dashboard permet de centraliser tout ça dans une interface unique, accessible à "
        "tous les membres du groupe sans avoir à toucher au code. Et au-delà de la visualisation, "
        "il joue un rôle analytique direct : avant de construire un modèle prédictif, "
        "il faut explorer les données et comprendre les tendances. "
        "C'est exactement ce que cet outil permet de faire."
    )

    # ── SLIDE 4 ───────────────────────────────────────────────────────────────
    h1("Slide 4 — Les 4 onglets")
    body(
        "Le dashboard est organisé en quatre onglets."
    )
    body(
        "Le premier, Vue d'ensemble, c'est la page d'accueil. Elle affiche les indicateurs "
        "clés — nombre d'échantillons, fibres détectées, contacts entre fibres, porosité moyenne — "
        "et un tableau qui compare tous les matériaux côte à côte."
    )
    body(
        "Le deuxième, Morphologie des fibres, analyse la distribution des diamètres. "
        "Un boxplot et une courbe de densité permettent de voir si un matériau a des fibres "
        "homogènes ou très dispersées, ce qui influence directement ses propriétés acoustiques "
        "selon Tran et al. (2024)."
    )
    body(
        "Le troisième, Propriétés acoustiques, montre les courbes d'absorption sonore "
        "de chaque échantillon sur cinq fréquences, et un scatter plot qui croise "
        "la porosité et la résistivité à l'air."
    )
    body(
        "Et le quatrième, Comparaison, c'est l'onglet central pour l'objectif du projet. "
        "Il compare les matériaux fréquence par fréquence, et met en relation le diamètre moyen "
        "des fibres et leur absorption à 1 kHz. C'est ce graphique qui permettra de tester "
        "directement l'hypothèse : est-ce que les fibres fines absorbent mieux le son ?"
    )

    # ── SLIDE 5 ───────────────────────────────────────────────────────────────
    h1("Slide 5 — Démonstration")
    body(
        "Je vous montre le dashboard maintenant. On va parcourir rapidement les quatre onglets "
        "pour voir ce que ça donne avec les données simulées."
    )
    note("(Ouvrir le navigateur sur http://127.0.0.1:8050 et parcourir les onglets)")

    # ── SLIDE 6 ───────────────────────────────────────────────────────────────
    h1("Slide 6 — Et après ?")
    body(
        "Pour la suite, la priorité est d'intégrer les vraies données dès qu'elles seront "
        "disponibles. L'outil est conçu pour ça : il suffit de remplacer les fichiers CSV, "
        "tout le reste fonctionne déjà. Et selon les besoins qui émergeront lors de l'analyse, "
        "de nouveaux graphiques pourront être ajoutés facilement."
    )
    body("Merci pour votre attention, je suis disponible pour les questions.")

    # ── SECTION GRAPHIQUES ────────────────────────────────────────────────────
    doc.add_page_break()

    p = doc.add_paragraph()
    r = p.add_run("Annexe — Description détaillée des graphiques")
    r.font.name = "Calibri"
    r.font.size = DPt(15)
    r.font.bold = True
    r.font.color.rgb = DRGBColor(0x1F, 0x49, 0x7D)

    note(
        "Cette section n'est pas à lire à l'oral. Elle sert à préparer les réponses aux "
        "questions des professeurs sur le fonctionnement des graphiques."
    )
    doc.add_paragraph()

    graph(
        "Tableau récapitulatif",
        "Vue d'ensemble",
        "À quoi il sert : donner une vue synthétique de tous les matériaux en un seul endroit, "
        "sans avoir à naviguer entre les onglets.\n\n"
        "Par quoi il est alimenté : les fichiers samples.csv, fibers.csv et acoustic_thermal.csv. "
        "Pour chaque matériau, les valeurs affichées sont des moyennes calculées sur l'ensemble "
        "des échantillons disponibles.\n\n"
        "Comment le lire : chaque ligne correspond à un matériau, chaque colonne à une métrique "
        "(diamètre moyen, longueur, porosité, orientation, absorption à 1 kHz). "
        "On compare directement les matériaux ligne par ligne.\n\n"
        "Comment l'exploiter : c'est le premier graphique à consulter. Il permet d'identifier "
        "rapidement quels matériaux se distinguent des autres et sur quels critères, "
        "avant d'approfondir dans les onglets suivants."
    )

    graph(
        "Boxplot des diamètres de fibres",
        "Morphologie des fibres",
        "À quoi il sert : comparer la distribution des diamètres entre matériaux et visualiser "
        "leur variabilité interne.\n\n"
        "Par quoi il est alimenté : le fichier fibers.csv, qui contient le diamètre de chaque "
        "fibre détectée pour chaque échantillon. Chaque boîte regroupe toutes les fibres "
        "d'un même matériau.\n\n"
        "Comment le lire : la ligne centrale est la médiane (valeur typique), la boîte encadre "
        "50 % des fibres (entre le 1er et le 3e quartile), et les traits montrent l'étendue "
        "des valeurs habituelles. Plus la boîte est haute, plus les fibres sont hétérogènes "
        "en taille — c'est la polydispersité.\n\n"
        "Comment l'exploiter : un matériau avec une boîte étroite a des fibres homogènes, "
        "ce qui simplifie la modélisation acoustique. Un matériau avec une boîte large présente "
        "une forte dispersion, qui peut influencer différemment les propriétés selon la fréquence "
        "sonore, comme le montrent Tran et al. (2024)."
    )

    graph(
        "Courbe de densité (KDE)",
        "Morphologie des fibres",
        "À quoi il sert : visualiser la forme précise de la distribution des diamètres, "
        "là où le boxplot ne montre que des quantiles.\n\n"
        "Par quoi il est alimenté : le fichier fibers.csv. Une estimation par noyau (KDE) est "
        "calculée sur les diamètres pour produire une courbe lissée.\n\n"
        "Comment le lire : l'axe horizontal représente le diamètre en micromètres, l'axe vertical "
        "la densité de probabilité. Un pic étroit et élevé signifie que la majorité des fibres "
        "ont un diamètre proche de cette valeur. Une courbe large et aplatie indique une grande "
        "diversité de tailles.\n\n"
        "Comment l'exploiter : en superposant les courbes de plusieurs matériaux, on voit si leurs "
        "distributions se chevauchent ou sont bien séparées. Ce graphique est directement inspiré "
        "de la Figure 3b de Tran et al. (2024), qui relie la forme de cette distribution aux "
        "propriétés acoustiques du matériau."
    )

    graph(
        "Courbes d'absorption acoustique",
        "Propriétés acoustiques",
        "À quoi il sert : montrer comment chaque échantillon absorbe le son en fonction de la "
        "fréquence, et comparer les profils d'absorption entre matériaux.\n\n"
        "Par quoi il est alimenté : le fichier acoustic_thermal.csv, qui contient pour chaque "
        "échantillon le coefficient d'absorption α mesuré à 5 fréquences : "
        "250 Hz, 500 Hz, 1 kHz, 2 kHz et 4 kHz.\n\n"
        "Comment le lire : l'axe vertical représente α, qui varie de 0 (le son rebondit "
        "totalement) à 1 (le son est totalement absorbé). Chaque courbe correspond à un "
        "échantillon. On teste plusieurs fréquences car un matériau peut bien absorber les "
        "graves et mal absorber les aigus, ou l'inverse.\n\n"
        "Comment l'exploiter : on cherche les matériaux dont la courbe reste élevée sur toutes "
        "les fréquences, ou au contraire ceux qui sont très performants sur une plage spécifique. "
        "Ces profils orienteront le choix des matériaux à modéliser en priorité."
    )

    graph(
        "Scatter porosité vs résistivité à l'air",
        "Propriétés acoustiques",
        "À quoi il sert : explorer la relation entre deux paramètres physiques clés qui influencent "
        "directement l'absorption acoustique.\n\n"
        "Par quoi il est alimenté : le fichier acoustic_thermal.csv, avec les colonnes porosity "
        "(proportion de vide dans le matériau) et airflow_resistivity (résistance que le matériau "
        "oppose au passage de l'air). Chaque point représente un échantillon.\n\n"
        "Comment le lire : un matériau très poreux (à droite) laisse l'air passer facilement et "
        "présente une faible résistivité. Un matériau dense (à gauche) résiste beaucoup au passage "
        "de l'air. Une courbe de tendance est ajustée automatiquement.\n\n"
        "Comment l'exploiter : les échantillons qui s'écartent fortement de la tendance sont "
        "particulièrement intéressants — ils peuvent indiquer une microstructure atypique ou une "
        "erreur de mesure à vérifier. Ce graphique aide aussi à calibrer les paramètres d'entrée "
        "du modèle acoustique."
    )

    graph(
        "Graphique à barres groupées",
        "Comparaison",
        "À quoi il sert : comparer directement tous les matériaux sur leur performance acoustique "
        "à chaque fréquence, en un seul graphique.\n\n"
        "Par quoi il est alimenté : le fichier acoustic_thermal.csv. Pour chaque matériau, on "
        "calcule la moyenne de α sur l'ensemble des échantillons, pour chacune des 5 fréquences.\n\n"
        "Comment le lire : chaque groupe de barres correspond à une fréquence. Au sein d'un groupe, "
        "chaque barre représente un matériau. Plus la barre est haute, meilleure est l'absorption.\n\n"
        "Comment l'exploiter : ce graphique permet de répondre directement à 'quel matériau est le "
        "plus performant ?' et de voir si ce classement change selon la fréquence. Un matériau qui "
        "domine sur toutes les fréquences est un candidat prioritaire pour la modélisation."
    )

    graph(
        "Scatter diamètre moyen vs absorption à 1 kHz",
        "Comparaison",
        "À quoi il sert : tester visuellement l'hypothèse centrale du projet — les fibres fines "
        "absorbent-elles mieux le son que les fibres épaisses ?\n\n"
        "Par quoi il est alimenté : une jointure entre fibers.csv (diamètre moyen par échantillon) "
        "et acoustic_thermal.csv (absorption à 1 kHz). Chaque point représente un échantillon.\n\n"
        "Comment le lire : l'axe horizontal représente le diamètre moyen des fibres en micromètres, "
        "l'axe vertical l'absorption à 1 kHz. 1 kHz correspond à la fréquence de la voix humaine, "
        "ce qui en fait la fréquence de référence la plus pertinente pour l'isolation acoustique.\n\n"
        "Comment l'exploiter : si les points forment une tendance descendante de gauche à droite, "
        "cela confirme que les fibres fines absorbent mieux le son. Si les points sont dispersés "
        "sans tendance claire, d'autres paramètres (porosité, longueur, orientation) jouent "
        "probablement un rôle important et devront être intégrés dans le modèle prédictif."
    )

    out = "docs/discours_oral_remy.docx"
    doc.save(out)
    print(f"DOCX discours créé : {out}")


if __name__ == "__main__":
    build_pptx()
    build_discours()
